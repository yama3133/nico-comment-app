"""
Lambda (Function URL) ハンドラ。
  action=presign  : 入力pptxアップロード用の署名付きPUT URLを発行
  action=generate : S3の入力pptxにコメント流れGIFを注入し、出力を署名付きGET URLで返す

GIF生成・pptx注入ロジックは CLI版 nico_cli.py を踏襲。フォントはコンテナ同梱のNoto Sans JP。
"""

import json
import os
import random
import re
import uuid

import boto3
from botocore.config import Config
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu

_s3_client = None


def s3():
    """S3クライアントを遅延生成（importだけでは認証に触れない）"""
    global _s3_client
    if _s3_client is None:
        _s3_client = boto3.client(
            "s3",
            region_name=os.environ.get("AWS_REGION", "ap-northeast-1"),
            config=Config(
                signature_version="s3v4",
                s3={"addressing_style": "virtual"},
            ),
        )
    return _s3_client


BUCKET = os.environ.get("BUCKET", "")
FONT_PATH = os.environ.get("FONT_PATH", "/var/task/fonts/NotoSansJP-VF.ttf")
URL_TTL = 24 * 3600  # 署名付きURLの有効期間(秒)

PALETTE = [
    (255, 255, 255), (255, 255, 0), (255, 136, 255), (136, 255, 255),
    (136, 255, 136), (255, 170, 68), (204, 153, 255),
]

CORS = {
    "Access-Control-Allow-Origin": "*",
    "Access-Control-Allow-Methods": "POST,OPTIONS",
    "Access-Control-Allow-Headers": "Content-Type",
}


def _resp(status, body):
    return {
        "statusCode": status,
        "headers": {"Content-Type": "application/json", **CORS},
        "body": json.dumps(body),
    }


def _hex_to_rgb(h):
    h = (h or "").lstrip("#")
    try:
        return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except (ValueError, IndexError):
        return None


def build_gif(comments, gif_path, *, width, height, font_size, fps,
              speed, interval, rows, seed=7):
    """コメントが右→左へ流れる透過アニメGIFを生成"""
    random.seed(seed)
    font = ImageFont.truetype(FONT_PATH, font_size)
    try:
        # variable fontをBold(700)に
        font.set_variation_by_axes([700])
    except Exception:  # noqa: BLE001
        pass
    tmp = Image.new("RGBA", (10, 10))
    td = ImageDraw.Draw(tmp)
    stroke = max(2, font_size // 10)

    items = []
    for i, (text, color) in enumerate(comments):
        if color is None:
            color = PALETTE[i % len(PALETTE)]
        bbox = td.textbbox((0, 0), text, font=font, stroke_width=stroke)
        tw = bbox[2] - bbox[0]
        row = i % rows
        y = (int(height * 0.03) + row * int(height * 0.94 / rows)
             + random.randint(-6, 6))
        items.append({"text": text, "color": color, "tw": tw,
                      "y": y, "start": i * interval})

    travel = (width + max(it["tw"] for it in items)) / speed
    total_time = items[-1]["start"] + travel
    n_frames = max(1, int(total_time * fps))

    p_frames = []
    for f in range(n_frames):
        t = f / fps
        img = Image.new("RGBA", (width, height), (0, 0, 0, 0))
        d = ImageDraw.Draw(img)
        for it in items:
            if t < it["start"]:
                continue
            x = width - speed * (t - it["start"])
            if x < -it["tw"]:
                continue
            d.text((x, it["y"]), it["text"], font=font,
                   fill=it["color"] + (255,), stroke_width=stroke,
                   stroke_fill=(0, 0, 0, 255))
        alpha = img.split()[3]
        p = img.convert("RGB").convert("P", palette=Image.ADAPTIVE, colors=255)
        mask = alpha.point(lambda a: 255 if a <= 128 else 0)
        p.paste(255, mask)
        p_frames.append(p)

    p_frames[0].save(
        gif_path, save_all=True, append_images=p_frames[1:],
        duration=int(1000 / fps), loop=0, transparency=255,
        disposal=2, optimize=False)


def inject(input_path, output_path, gif_path, target_slides):
    """指定スライド(1始まり)にGIFを全面で重ねる"""
    prs = Presentation(input_path)
    w, h = prs.slide_width, prs.slide_height
    slides = list(prs.slides)
    total = len(slides)
    for n in target_slides:
        if 1 <= n <= total:
            slides[n - 1].shapes.add_picture(gif_path, Emu(0), Emu(0),
                                             width=w, height=h)
    prs.save(output_path)
    return w, h


def handle_presign(body):
    key = f"uploads/{uuid.uuid4().hex}.pptx"
    url = s3().generate_presigned_url(
        "put_object",
        Params={
            "Bucket": BUCKET, "Key": key,
            "ContentType": "application/vnd.openxmlformats-officedocument."
                           "presentationml.presentation",
        },
        ExpiresIn=3600,
    )
    return _resp(200, {"uploadUrl": url, "inputKey": key})


def handle_generate(body):
    input_key = body["inputKey"]
    slides = [int(n) for n in body.get("slides", [])]
    raw_comments = body.get("comments", [])
    s = body.get("settings", {})
    speed = float(s.get("speed", 220))
    interval = float(s.get("interval", 0.42))
    rows = int(s.get("rows", 6))
    font_size = int(s.get("fontSize", 30))

    comments = []
    for c in raw_comments:
        text = (c.get("text") or "").strip()
        if not text:
            continue
        comments.append((text, _hex_to_rgb(c.get("color"))))
    if not comments:
        return _resp(400, {"error": "コメントがありません"})
    if not slides:
        return _resp(400, {"error": "対象スライドがありません"})

    in_path = f"/tmp/{uuid.uuid4().hex}.pptx"
    gif_path = f"/tmp/{uuid.uuid4().hex}.gif"
    out_path = f"/tmp/{uuid.uuid4().hex}.pptx"

    s3().download_file(BUCKET, input_key, in_path)

    # スライドサイズ比率に合わせてGIF高さを決定
    prs = Presentation(in_path)
    ratio = prs.slide_height / prs.slide_width
    gif_w = 960
    gif_h = int(gif_w * ratio)
    del prs

    build_gif(comments, gif_path, width=gif_w, height=gif_h,
              font_size=font_size, fps=12, speed=speed,
              interval=interval, rows=rows)
    inject(in_path, out_path, gif_path, slides)

    out_key = f"outputs/{uuid.uuid4().hex}.pptx"
    s3().upload_file(
        out_path, BUCKET, out_key,
        ExtraArgs={
            "ContentType": "application/vnd.openxmlformats-officedocument."
                           "presentationml.presentation",
            "ContentDisposition": 'attachment; filename="commented.pptx"',
        },
    )
    url = s3().generate_presigned_url(
        "get_object",
        Params={"Bucket": BUCKET, "Key": out_key},
        ExpiresIn=URL_TTL,
    )
    return _resp(200, {"downloadUrl": url, "outputKey": out_key})


def handler(event, context):
    method = (event.get("requestContext", {})
              .get("http", {}).get("method", "POST"))
    if method == "OPTIONS":
        return {"statusCode": 204, "headers": CORS, "body": ""}

    qs = event.get("queryStringParameters") or {}
    action = qs.get("action", "")
    try:
        body = json.loads(event.get("body") or "{}")
    except json.JSONDecodeError:
        return _resp(400, {"error": "不正なJSON"})

    try:
        if action == "presign":
            return handle_presign(body)
        if action == "generate":
            return handle_generate(body)
        return _resp(400, {"error": f"未知のaction: {action}"})
    except Exception as e:  # noqa: BLE001
        return _resp(500, {"error": str(e)})
